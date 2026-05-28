import random
from math import gcd
from fractions import Fraction
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re
from fractions import Fraction
from sympy import symbols, Eq, solve, isprime
import json
import os
import sys
import uuid
import math
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

x = symbols('x')

# levels_to_generate = [1,2,3,4,5,6,7,8,9]
levels_to_generate = [6]
SHAPE_LABEL_FONT_SIZE = 28   # Font size for dimension/angle labels on all shapes
SHAPE_LABEL_FONT_WEIGHT = 'bold'  # Font weight for dimension/angle labels on all shapes
# Simple color name to RGB map
COLOR_MAP = {
    "black": RGBColor(0, 0, 0),
    "blue": RGBColor(0, 0, 255),
    "red": RGBColor(255, 0, 0),
    "green": RGBColor(0, 128, 0),
    "dark_gray": RGBColor(100, 100, 100),
}

QUESTION_TYPE_FILE = "question_data.json"

def load_question_types():
    if os.path.exists(QUESTION_TYPE_FILE):
        with open(QUESTION_TYPE_FILE, "r") as f:
            return json.load(f)
    return []

question_typeX = load_question_types()

question_type = {
    1: {
        "add": [
            {"qty":30, "min": [1,0], "max": [6,6], "tiers": [1]},
            {"qty":30, "min": [1,4], "max": [6,10], "tiers": [2]},
            {"qty":10, "min": [1,10], "max": [7,13], "tiers": [3]},
        ],
        "subtract": [
            {"qty":10, "min": [15,1], "max": [20,8], "tiers": [3,4]},
            ],
        "multiply": [{"qty":10, "min": 1, "max": 8, "tiers": [4,5]}],
        "place_value": [{"qty":10, "min": 1, "max": 99, "tiers": [3,4], "fontsize": 80}],
        "place_value_reverse": [{"qty":20, "min": 10, "max": 99, "tiers": [3,4,5], "fontsize": 80}],
    },
    2: {
        "add": [
            {"qty":30, "min": 1, "max": 9, "tiers": [1]},
            {"qty":15, "min": [2,10], "max": [9,20], "tiers": [2]},
            {"qty":30, "min": 10, "max": 30, "tiers": [3,4]},
        ],
        "add3": [
            {"qty":15, "min": 1, "max": 10, "tiers": [4,5]},],
        "subtract": [
            {"qty":15, "min": 1, "max": 20, "tiers": [2]},
            {"qty":20, "min": 5, "max": 30, "tiers": [3,4]}
        ],
        "multiply": [{"qty":20, "min": 1, "max": 10, "tiers": [3,4,5]}],
        "place_value": [{"qty":10, "min": 1, "max": 999, "tiers": [3,4], "fontsize": 80}],
        "place_value_reverse": [{"qty":10, "min": 10, "max": 9999, "tiers": [3,4,5], "fontsize": 80}],
    },
    3: {
        "add": [
            {"qty":15, "min": 5, "max": 10, "tiers": [1]},
            {"qty":10, "min": 10, "max": 50, "tiers": [2]},
            {"qty":10, "min": [2,15], "max": [9,30], "tiers": [2]},
            ],
        "subtract": [
            {"qty":15, "min": 5, "max": 15, "tiers": [1]},
            {"qty":10, "min": 10, "max": 20, "tiers": [2]}
            ],
        "add3": [
            {"qty":20, "min": 1, "max": 20, "tiers": [3,4]}],
        "multiply": [
            {"qty":15, "min": 2, "max": 5, "tiers": [3]},
            {"qty":20, "min": 4, "max": 10, "tiers": [4,5]},
            ],
        "divide": [
            {"qty":20, "min": 1, "max": 20, "tiers": [3,4]}],
        "area": [
            {"qty":10, "min": 3, "max": 10, "tiers": [3,4]}],
        "place_value_reverse": [
            {"qty":20, "min": 100, "max": 9999, "tiers": [4,5], "fontsize": 80}],
    },
    4: {
        "add": [
            {"qty":15, "min": [2,10], "max": [9,15], "tiers": [1]},
            {"qty":15, "min": 5, "max": 15, "tiers": [1]},
            {"qty":15, "min": 10, "max": 20, "tiers": [2]},
            ],
        "add3": [
            {"qty":15, "min": 2, "max": 10, "tiers": [2]}],
        "add_fraction_same_denominator": [
            {"qty":15, "min": 2, "max": 10, "tiers": [3,4]}],
        "add_fraction_different_denominator": [
            {"qty":15, "min": 2, "max": 6, "tiers": [4,5]}],
        "subtract": [
            {"qty":15, "min": 5, "max": 20, "tiers": [1]},
            {"qty":15, "min": 10, "max": 40, "tiers": [3,4]},
            ],
        "multiply": [
            {"qty":15, "min": 3, "max": 7, "tiers": [3]},
            {"qty":15, "min": 5, "max": 10, "tiers": [3,4]},
            ],
        "divide": [
            {"qty":20, "min": 2, "max": 100, "tiers": [3,4,5]}],
        "area": [
            {"qty":10, "min": 4, "max": 12, "tiers": [3,4,5]}],
        "place_value_reverse": [
            {"qty":20, "min": 10, "max": 9999, "tiers": [3,4], "fontsize": 70}],
    },
    5: {
        "add": [
            {"qty":30, "min": [2,9], "max": [9,20], "tiers": [1]},
            {"qty":15, "min": 10, "max": 30, "tiers": [2]},
            {"qty":10, "min": 15, "max": 50, "tiers": [3,4,5]},
            ],
        "add3": [{"qty":10, "min": 5, "max": 15, "tiers": [3,4]}],
        "subtract": [
            {"qty":15, "min": 5, "max": 20, "tiers": [2]},
            {"qty":15, "min": 15, "max": 40, "tiers": [3,4]},
            ],
        "multiply": [
            {"qty":15, "min": 3, "max": 10, "tiers": [3]},
            ],
        "divide": [
            {"qty":15, "min": 2, "max": 100, "tiers": [4,5]}],
        "place_value_reverse": [
            {"qty":15, "min": 100, "max": 9999, "tiers": [3,4,5], "fontsize": 70}],
        "add_dec1": [
            {"qty":10, "min": 1, "max": 10, "tiers": [4,5]}],
        "add_fraction_same_denominator": [
            {"qty":15, "min": 2, "max": 10, "tiers": [3,4]}],
        "add_fraction_different_denominator": [
            {"qty":15, "min": 2, "max": 7, "tiers": [4,5]}],

    },
    6: {
        "add": [
            {"qty":15, "min": 10, "max": 30, "tiers": [1]},
            {"qty":15, "min": 20, "max": 60, "tiers": [2]},
            ],
        "add3": [{"qty":10, "min": 10, "max": 40, "tiers": [3,4]}],
        "perc10": [
            {"qty":10, "min": 10, "max": 500, "tiers": [3,4]}
            ],
        "add_dec1": [{"qty":10, "min": 1, "max": 10, "tiers": [4,5]}],
        "subtract": [
            {"qty":15, "min": 5, "max": 15, "tiers": [1]},
            {"qty":15, "min": 10, "max": 20, "tiers": [2]},
            {"qty":15, "min": 15, "max": 50, "tiers": [3,4]},
            ],
        "multiply": [
            {"qty":30, "min": 3, "max": 10, "tiers": [3]},
            ],
        "divide": [
            {"qty":15, "min": 2, "max": 40, "tiers": [4,5]}],
        "area": [
            {"qty":10, "min": 5, "max": 15, "tiers": [4,5]}],
        "missing_angle_triangle": [
            {"qty":8, "min": 1, "max": 10, "tiers": [3,4,5]}],
        "place_value_reverse": [
            {"qty":15, "min": 1000, "max": 99999, "tiers": [3,4,5], "fontsize": 70}],
        "add_fraction_same_denominator": [
            {"qty":10, "min": 2, "max": 10, "tiers": [3]}],
        "add_fraction_different_denominator": [
            {"qty":10, "min": 2, "max": 9, "tiers": [4]}],

    },
    7: {
        "add": [
            {"qty":15, "min": [5,10], "max": [10,20], "tiers": [1]},
            {"qty":10, "min": [10,15], "max": [20,30], "tiers": [2]},
            ],
        "add3": [{"qty":20, "min": [3,10,15], "max": [10,20,30], "tiers": [3,4]}],
        "perc10": [
            {"qty":10, "min": 10, "max": 500, "tiers": [3,4]}
            ],
        "add_dec1": [
            {"qty":10, "min": 1, "max": 20, "tiers": [4,5]}],
        "subtract": [
            {"qty":15, "min": 5, "max": 20, "tiers": [1]},
            {"qty":10, "min": 10, "max": 30, "tiers": [2]},
            ],
        "multiply": [
            {"qty":10, "min": 6, "max": 10, "tiers": [2]},
            ],
        "divide<11": [
            {"qty":10, "min": 2, "max": 80, "tiers": [4,5]}],
        "place_value_reverse": [
            {"qty":15, "min": 1000, "max": 99999, "tiers": [3,4,5], "fontsize": 70}],
        "missing_angle_triangle": [
            {"qty":8, "min": 1, "max": 10, "tiers": [3,4,5]}],
        "add_fraction_same_denominator": [
            {"qty":10, "min": 2, "max": 9, "tiers": [3]}],
        "add_fraction_different_denominator": [
            {"qty":20, "min": 2, "max": 9, "tiers": [4,5]}],
        "linear_equation_a*x=c": [
            {"qty":10, "min": 1, "max": 10, "tiers": [3,4], "fontsize": 90},
            ],
        "linear_equation_a(x+b)=c": [
            {"qty":20, "min": 1, "max": 10, "tiers": [4,5], "fontsize": 90},
            ],
        "linear_equation_ax+b=c": [
            {"qty":20, "min": 1, "max": 10, "tiers": [4,5], "fontsize": 90},
            ],
        "linear_equation_x+a=c": [
            {"qty":10, "min": 1, "max": 10, "tiers": [4], "fontsize": 90},
            ],
        "linear_equation_x-a=c": [
            {"qty":10, "min": 1, "max": 10, "tiers": [4], "fontsize": 90},
            ],

    },
    8: {
        "add": [
            {"qty":15, "min": 6, "max": 20, "tiers": [1]},
            {"qty":10, "min": 10, "max": 30, "tiers": [2]},
            ],
        "add3": [{"qty":20, "min": [3,10,15], "max": [12,20,50], "tiers": [3,4]}],
        "perc10": [
            {"qty":10, "min": 10, "max": 500, "tiers": [3,4]}
            ],
        "add_dec1": [{"qty":10, "min": 1, "max": 20, "tiers": [4,5]}],
        "subtract": [
            {"qty":15, "min": 5, "max": 30, "tiers": [1]},
            {"qty":10, "min": 10, "max": 50, "tiers": [2,3,4]},
            ],
        "linear_equation_a*x=c": [
            {"qty":15, "min": -10, "max": 10, "tiers": [3], "fontsize": 90},
            ],
        "linear_equation_a(x+b)=c": [
            {"qty":20, "min": 1, "max": 10, "tiers": [4,5], "fontsize": 90},
            ],
        "linear_equation_ax+b=c": [
            {"qty":20, "min": 1, "max": 10, "tiers": [3,4,5], "fontsize": 90},
            ],
        "linear_equation_x+a=c": [
            {"qty":10, "min": 1, "max": 10, "tiers": [4], "fontsize": 90},
            ],
        "linear_equation_x-a=c": [
            {"qty":10, "min": 1, "max": 10, "tiers": [4], "fontsize": 90},
            ],
        "multiply": [
            {"qty":10, "min": 6, "max": 10, "tiers": [3,4]},
            ],
        "divide<11": [
            {"qty":10, "min": 2, "max": 80, "tiers": [4,5]}],
        "place_value_reverse": [
            {"qty":10, "min": 1000, "max": 99999, "tiers": [3,4,5], "fontsize": 70}],
        "add_fraction_same_denominator": [
            {"qty":10, "min": 2, "max": 9, "tiers": [2]},
        ],
        "add_fraction_different_denominator": [
            {"qty":10, "min": 2, "max": 9, "tiers": [4,5]}],
    },
    9: {
        "add": [
            {"qty":15, "min": 10, "max": 20, "tiers": [1]},
            {"qty":15, "min": [-10,1], "max": [-1,10], "tiers": [2]},
            {"qty":10, "min": [-10,-10], "max": [-1,10], "tiers": [3,4]},
            ],
        "add3": [{"qty":20, "min": [3,10,15], "max": [20,40,50], "tiers": [3,4]}],
        "perc10": [
            {"qty":10, "min": 10, "max": 1000, "tiers": [3,4,5]}],
        "add_dec1": [{"qty":20, "min": 1, "max": 20, "tiers": [4,5]}],
        "subtract": [
            {"qty":15, "min": [10,5], "max": [20,15], "tiers": [1]},
            {"qty":15, "min": [15,7], "max": [40,17], "tiers": [2]},
            {"qty":10, "min": [20,7], "max": [60,30], "tiers": [3,4,5]},
            {"qty":10, "min": [10,-20], "max": [50,-2], "tiers": [3,4,5]},
            ],
        "linear_equation_a*x=c": [
            {"qty":10, "min": -10, "max": 10, "tiers": [3], "fontsize": 90},
            ],
        "linear_equation_a(x+b)=c": [
            {"qty":20, "min": -10, "max": 10, "tiers": [4,5], "fontsize": 90},
            ],
        "linear_equation_ax+b=c": [
            {"qty":20, "min": -10, "max": 10, "tiers": [4,5], "fontsize": 90},
            ],
        "linear_equation_x+a=c": [
            {"qty":10, "min": -10, "max": 10, "tiers": [4], "fontsize": 90},
            ],
        "linear_equation_x-a=c": [
            {"qty":10, "min": -10, "max": 10, "tiers": [4], "fontsize": 90},
            ],
        "multiply": [
            {"qty":10, "min": 6, "max": 10, "tiers": [3]},
            ],
        "divide<11": [
            {"qty":10, "min": 2, "max": 80, "tiers": [4,5]}],
        "place_value_reverse": [
            {"qty":15, "min": 1000, "max": 999999, "tiers": [3,4,5], "fontsize": 70}],
        "add_fraction_different_denominator": [
            {"qty":20, "min": 2, "max": 9, "tiers": [3,4,5]}],
    },
    10: {
        "add": [
            {"qty":15, "min": 10, "max": 30, "tiers": [1]},
            {"qty":15, "min": [-10,1], "max": [-1,10], "tiers": [2]},
            {"qty":10, "min": [-10,-10], "max": [-1,-1], "tiers": [3,4]},
            ],
        "add3": [
            {"qty":10, "min": 10, "max": 70, "tiers": [3,4]}],
        "perc10": [
            {"qty":10, "min": 10, "max": 1000, "tiers": [3,4,5]}],
        "add_dec1": [{"qty":20, "min": 1, "max": 20, "tiers": [4,5]}],
        "subtract": [
            {"qty":15, "min": 5, "max": 20, "tiers": [1]},
            {"qty":15, "min": 7, "max": 40, "tiers": [2]},
            {"qty":10, "min": 7, "max": 60, "tiers": [3,4,5]},
            {"qty":10, "min": [10,-20], "max": [50,-2], "tiers": [3,4,5]},
            ],
        "linear_equation_a*x=c": [
            {"qty":10, "min": -10, "max": 10, "tiers": [3], "fontsize": 90},
            ],
        "linear_equation_a(x+b)=c": [
            {"qty":10, "min": -10, "max": 10, "tiers": [4], "fontsize": 90},
            ],
        "linear_equation_ax+b=c": [
            {"qty":10, "min": -10, "max": 10, "tiers": [4], "fontsize": 90},
            ],
        "multiply": [
            {"qty":10, "min": 6, "max": 10, "tiers": [3]},
            ],
        "divide<11": [
            {"qty":10, "min": 2, "max": 80, "tiers": [4,5]}],
        "place_value_reverse": [
            {"qty":5, "min": 1000, "max": 999999, "tiers": [3,4,5], "fontsize": 70}],
        "add_fraction_different_denominator": [
            {"qty":10, "min": 2, "max": 9, "tiers": [3,4,5]}],
    },

}

def unicode_fraction(numerator, denominator):
    super_digits = {
        '0': '⁰', '1': '¹', '2': '²', '3': '³', '4': '⁴', '5': '⁵', '6': '⁶', '7': '⁷', '8': '⁸', '9': '⁹', '-': '⁻'
    }
    sub_digits = {
        '0': '₀', '1': '₁', '2': '₂', '3': '₃', '4': '₄', '5': '₅', '6': '₆', '7': '₇', '8': '₈', '9': '₉', '-': '₋'
    }
    def to_super(num):
        return ''.join(super_digits.get(ch, ch) for ch in str(num))

    def to_sub(num):
        return ''.join(sub_digits.get(ch, ch) for ch in str(num))
    
    return f"{to_super(numerator)}⁄{to_sub(denominator)}"  # Using Unicode fraction slash (U+2044)

def add_textbox(slide, text, left, top, width, height, default_font_size=44, 
                default_color = COLOR_MAP['black'], bold=True, 
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):

    def parse_size(size_str, default_size_pt):
        if not size_str:
            return default_size_pt
        if size_str.endswith('%'):
            percent = float(size_str[:-1])
            return default_size_pt * percent / 100
        else:
            return float(size_str)

    def add_run(paragraph, text, size, color, italic=False, font_name=None, bold=False):
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        if font_name:
            run.font.name = font_name        
        return run

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True  # This is default, but explicitly set it
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.line_spacing = 0.9  # proportion, not points

    pattern = re.compile(
        r'<font'                                               # Start of tag
        r'(?:\s+size=([\d]+%?))?'                              # Group 1: size
        r'(?:\s+color=(\w+))?'                                 # Group 2: color
        r'(?:\s+italic=(true|false))?'                         # Group 3: italic
        r'(?:\s+font=(?:"([^"]+)"|\'([^\']+)\'|(\S+)))?'       # Groups 4–6: font name
        r'\s*>(.*?)</font>',                                   # Group 7: inner text
        re.DOTALL
    )

    
    matches = list(pattern.finditer(text))
    last_end = 0

    for match in matches:
        if match.start() > last_end:
            add_run(p, text[last_end:match.start()], default_font_size, default_color)

        size_str = match.group(1)
        color_str = match.group(2)
        italic_str = match.group(3)
        font_name = match.group(4) or match.group(5) or match.group(6)
        content = match.group(7)

        size = parse_size(size_str, default_font_size)
        color = COLOR_MAP.get(color_str.lower(), default_color) if color_str else default_color
        italic = italic_str.lower() == "true" if italic_str else False

        add_run(p, content, size, color, italic=italic, font_name=font_name)
        last_end = match.end()

    # Remaining plain text after last match
    if last_end < len(text):
        add_run(p, text[last_end:], default_font_size, default_color)

    p.alignment = align
    text_frame.vertical_anchor = valign

    return textbox


def draw_area_shape(slide, shape_info, left, top, display_width, display_height):
    """Draw a rectangle or triangle and label its dimensions."""
    label_font_size = SHAPE_LABEL_FONT_SIZE

    width_value = shape_info.get("width") or shape_info.get("base")
    height_value = shape_info["height"]
    max_value = max(width_value, height_value, 1)

    shape_width = display_width * (width_value / max_value)
    shape_height = display_height * (height_value / max_value)
    shape_left = left + (display_width - shape_width) / 2
    shape_top = top + (display_height - shape_height) / 2

    if shape_info["shape"] == "rectangle":
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, shape_left, shape_top, shape_width, shape_height)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, shape_left, shape_top, shape_width, shape_height)

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 255)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(2)

    if shape_info["shape"] == "rectangle":
        add_textbox(slide, f"{width_value} cm", shape_left, shape_top + shape_height + Inches(0.1), shape_width, Inches(0.4), default_font_size=label_font_size, align=PP_ALIGN.CENTER)
        add_textbox(slide, f"{height_value} cm", shape_left - Inches(0.9), shape_top, Inches(0.9), shape_height, default_font_size=label_font_size, align=PP_ALIGN.CENTER)
    else:
        add_textbox(slide, f"{shape_info['base']} cm", shape_left, shape_top + shape_height + Inches(0.1), shape_width, Inches(0.4), default_font_size=label_font_size, align=PP_ALIGN.CENTER)
        add_textbox(slide, f"{height_value} cm", shape_left - Inches(0.9), shape_top, Inches(0.9), shape_height, default_font_size=label_font_size, align=PP_ALIGN.CENTER)


def draw_triangle_angle(slide, qinfo, left, top, display_width, display_height, show_answer=False):
    """Render a triangle using matplotlib, annotate angles and insert into slide.

    qinfo: {"angles": [a1,a2,a3], "missing_index": i}
    """
    angles = qinfo.get("angles", [60,60,60])
    missing_index = qinfo.get("missing_index", 2)

    # Convert to radians
    a_rad = math.radians(angles[0])
    b_rad = math.radians(angles[1])
    c_rad = math.radians(angles[2])

    # Use base AB = 1.0
    c_len = 1.0
    # By law of sines: a/sin(A)=b/sin(B)=c/sin(C)=scale
    scale = c_len / math.sin(c_rad)
    a_len = math.sin(a_rad) * scale  # side opposite A (BC)
    b_len = math.sin(b_rad) * scale  # side opposite B (AC)

    # Coordinates: A at (0,0), B at (c_len,0). Find C such that AC=b_len, BC=a_len
    # x coordinate from A: x = (b^2 + c^2 - a^2)/(2c)
    x = (b_len*b_len + c_len*c_len - a_len*a_len) / (2*c_len)
    y_sq = max(0.0, b_len*b_len - x*x)
    y = math.sqrt(y_sq)

    # Build polygon coordinates
    A = (0.0, 0.0)
    B = (c_len, 0.0)
    C = (x, y)
    xs = [A[0], B[0], C[0], A[0]]
    ys = [A[1], B[1], C[1], A[1]]

    # Convert display sizes (pptx Inches) to float inches if possible
    try:
        dw = float(display_width.inches)
    except Exception:
        try:
            dw = float(display_width)
        except Exception:
            dw = 5.0
    try:
        dh = float(display_height.inches)
    except Exception:
        try:
            dh = float(display_height)
        except Exception:
            dh = 3.0

    # Limit figure size to reasonable bounds to avoid huge image buffers
    fig_w = max(3.0, min(dw, 6.0))
    fig_h = max(3.0, min(dh, 4.0))
    fig = plt.figure(figsize=(fig_w, fig_h), dpi=100)
    ax = fig.add_subplot(111)
    ax.fill(xs, ys, color="#F0F0FF", edgecolor="black")
    ax.plot(xs, ys, color="black")

    # Annotate angles — offset each label away from the triangle centroid
    centroid_x = (A[0] + B[0] + C[0]) / 3
    centroid_y = (A[1] + B[1] + C[1]) / 3
    offset_dist = 0.12  # distance to push label away from vertex, in data units
    for idx, (vx, vy) in enumerate([A, B, C]):
        if idx == missing_index and not show_answer:
            label = "?"
        else:
            label = f"{angles[idx]}°"
        # Direction from centroid to vertex (outward)
        dx = vx - centroid_x
        dy = vy - centroid_y
        length = math.sqrt(dx * dx + dy * dy) or 1.0
        ox = (dx / length) * offset_dist
        oy = (dy / length) * offset_dist
        ax.text(vx + ox, vy + oy, label, fontsize=SHAPE_LABEL_FONT_SIZE, weight=SHAPE_LABEL_FONT_WEIGHT,
                ha='center', va='center')

    ax.set_aspect('equal')
    ax.axis('off')

    img_dir = os.path.join(os.getcwd(), 'temp_images')
    os.makedirs(img_dir, exist_ok=True)
    img_path = os.path.join(img_dir, f"triangle_{uuid.uuid4().hex}.png")
    fig.savefig(img_path, bbox_inches='tight', dpi=150, transparent=True)
    plt.close(fig)

    # Fit image within display_width x display_height, preserving aspect ratio
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        img_w, img_h = im.size

    img_aspect = img_w / img_h
    fitted_w = display_width
    fitted_h = display_width / img_aspect
    if fitted_h > display_height:
        fitted_h = display_height
        fitted_w = display_height * img_aspect

    offset_x = left + (display_width - fitted_w) / 2
    offset_y = top + (display_height - fitted_h) / 2

    # Insert image into slide
    try:
        slide.shapes.add_picture(img_path, offset_x, offset_y, width=fitted_w, height=fitted_h)
    except Exception:
        slide.shapes.add_picture(img_path, left, top)
    return img_path

def format_question_summary(question):
    if isinstance(question, dict) and question.get("type") == "area":
        if question["shape"] == "rectangle":
            return f"Rectangle area with width {question['width']} cm and height {question['height']} cm"
        return f"Triangle area with base {question['base']} cm and height {question['height']} cm"
    if isinstance(question, dict) and question.get("type") == "triangle_angle":
        missing = question.get('missing_index', 2)
        return f"Triangle with one missing angle (angle {missing+1} to find)"
    return question


def generate_number(min_val, max_val, tier, tier_range):
    """Generate a number scaled by tier difficulty, including min_val and max_val."""
    # Total number of tiers
    num_tiers = len(tier_range)

    # We can pass in list of min_val and max_val for some operations so need to handle that. 
    # The operation will pull the value it wants and pass it to this function.
    if isinstance(min_val, list):
        min_val = min_val[0]
        
    if isinstance(max_val, list):
        max_val = max_val[0]
        
    # Size of each tier range
    range_size = (max_val - min_val + 1) / num_tiers  # +1 to include max_val

    # Determine this tier's range
    tier_index = tier - min(tier_range)
    tier_min = int(min_val + tier_index * range_size)
    tier_max = int(min_val + (tier_index + 1) * range_size - 1)

    # Clamp to global min and max to avoid overshooting
    tier_min = int(max(min_val, tier_min))
    tier_max = int(min(max_val, tier_max))

    return random.randint(tier_min, tier_max)

def simplify_fraction(num, denom):
    """Simplify a fraction and return as a string."""
    if denom == 0:
        return "undefined"
    g = gcd(num, denom)
    num, denom = num // g, denom // g
    if denom == 1:
        return str(num)
    if denom < 0:
        num, denom = -num, -denom
    return f"{num}/{denom}"

def generate_question(op_type, min_val, max_val, tier, tier_range):
    """Generate a question and answer based on operation type."""
    a = generate_number(min_val, max_val, tier, tier_range)
    b = generate_number(min_val, max_val, tier, tier_range)

    def pretty_fraction(frac):
        """Return pretty version of a fraction or mixed number."""
        if frac.denominator == 1:
            return str(frac.numerator)
        mixed = divmod(frac.numerator, frac.denominator)
        pretty = unicode_fraction(mixed[1], frac.denominator)

        if mixed[0] > 0:
            return f"{mixed[0]}{pretty}"  
        else:
            return pretty
    
    def linear_equation(op_type, min_val, max_val, tier, tier_range):
        a = generate_number(min_val, max_val, tier, tier_range)
        b = generate_number(min_val, max_val, tier, tier_range)
        xa = generate_number(min_val, max_val, tier, tier_range)  # known solution
        x = "<font size=120% italic=true font='Times New Roman'>x</font>"  # x with a subscript

        form = op_type.split("_equation_")[-1]

        if form == 'ax+b=c':
            c = a * xa + b
            equation = f"{a}{x} + {b} = {c}"

        elif form == 'a(x+b)=c':
            c = a * (xa + b)
            equation = f"{a}({x} + {b}) = {c}"

        elif form == 'x+a=c':
            c = xa + a
            equation = f"{x} + {a} = {c}"

        elif form == 'x-a=c':
            c = xa - a
            equation = f"{x} - {a} = {c}"

        elif form == 'a*x=c':
            c = a * xa
            equation = f"{a}{x} = {c}"

        return equation, xa

    if op_type == "add":

        if isinstance(min_val, list) and isinstance(max_val, list):
            b = generate_number(min_val[1], max_val[1], tier, tier_range)

        question = f"{a} + {b}"
        answer = a + b
        return question, f"{answer:,}"
    
    elif op_type.startswith("linear_equation"):
        question, answer = linear_equation(op_type, min_val, max_val, tier, tier_range)
        return question, f"{answer:,}"
    
    elif op_type == "add3":  # New operation for adding three numbers
        if isinstance(min_val, list) and isinstance(max_val, list):
            b = generate_number(min_val[1], max_val[1], tier, tier_range)
            c = generate_number(min_val[2], max_val[2], tier, tier_range)
        else:
            c = generate_number(min_val, max_val, tier, tier_range)
        question = f"{a} + {b} + {c}"
        answer = a + b + c
        return question, f"{answer:,}"
      
    elif op_type == "area":
        shape = random.choice(["rectangle", "triangle"])
        if shape == "rectangle":
            width_val = generate_number(min_val, max_val, tier, tier_range)
            height_val = generate_number(min_val, max_val, tier, tier_range)
            question = {
                "type": "area",
                "shape": "rectangle",
                "width": width_val,
                "height": height_val,
            }
            answer = width_val * height_val
            return question, f"{answer} cm²"

        base = generate_number(min_val, max_val, tier, tier_range)
        height_val = generate_number(min_val, max_val, tier, tier_range)
        while (base * height_val) % 2 != 0:
            height_val = generate_number(min_val, max_val, tier, tier_range)
        area = (base * height_val) // 2
        question = {
            "type": "area",
            "shape": "triangle",
            "base": base,
            "height": height_val,
        }
        return question, f"{area} cm²"

    # --- Missing angle in a triangle ---
    elif op_type == "missing_angle_triangle":
        # Generate two angles such that the third is positive
        while True:
            a1 = random.randint(20, 80)
            a2 = random.randint(20, 80)
            if a1 + a2 < 170:
                break
        a3 = 180 - (a1 + a2)
        angles = [a1, a2, a3]
        missing_index = random.choice([0, 1, 2])
        q = {
            "type": "triangle_angle",
            "angles": angles,
            "missing_index": missing_index,
        }
        answer = f"{angles[missing_index]}°"
        return q, answer

    elif op_type == "add_dec1":  # New operation for adding decimals
        a = generate_number(min_val*10, max_val*10, tier, tier_range) / 10
        b = generate_number(min_val*10, max_val*10, tier, tier_range) / 10
        question = f"{a:.1f} + {b:.1f}"
        answer = a + b
        return question, f"{answer:,.1f}"

    elif op_type == "perc10":  # New operation for adding decimals
        p = random.randint(1, 9) * 10
        b = int(b / 10) * 10  # Ensure b is a multiple of 10
        question = f"{p}% of {b:,}"
        answer = b * (p / 100)
        return question, f"{int(answer):,}" if answer == int(answer) else f"{answer:,}"

    elif op_type == "subtract":
        if isinstance(min_val, list) and isinstance(max_val, list): # If we are controlling the range precisely then we probably dont care if it yields a negative result
            b = generate_number(min_val[1], max_val[1], tier, tier_range)
        elif a < b:      # Ensure non-negative result by making a >= b
            a, b = b, a  # Swap if a < b
        if a == b:
            a += 2
        question = f"{a} - {b}"
        answer = a - b
        return question, f"{answer:,}"

    elif op_type == "multiply":
        question = f"{a} × {b}"
        answer = a * b
        return question, f"{answer:,}"

    elif op_type == "divide<11":
        d = random.randint(2, 10)  # Denominator between 1 and 10
        b = generate_number(d, max_val, tier, tier_range)
        a = int(b/d)

        product = a * d
        question = f"{product} ÷ {d}"
        answer = a
        return question, f"{answer:,}"

    elif op_type == "divide":
        a = generate_number(min_val, max_val / 2, tier, tier_range)
        if a == 0:
            a = 1
        new_max_val = int(max_val // a)  # work out the max product to still fit between our range
        if new_max_val <3:
            new_max_val = 3
        b = generate_number(2, new_max_val, tier, tier_range)
        if(isprime(b)):
            b += 1  # Avoid prime numbers to ensure integer division
        # Ensure integer division result
        product = a * b
        question = f"{product} ÷ {b}"
        answer = a
        return question, f"{answer:,}"

    elif op_type == "add_fraction_same_denominator":
        denoms = range(min(a,b), max(a,b)+1)
        if len(denoms) == 1:
            denom = denoms[0]
        else:
            denom = random.choice(denoms)
        a_num = random.randint(1, denom - 1)
        # b_num = random.randint(1, denom - a_num)  # ensure < 1
        b_num = random.randint(1, denom )
        a = Fraction(a_num, denom)
        b = Fraction(b_num, denom)

        answer = a + b
        question = f"{unicode_fraction(a_num, denom)} + {unicode_fraction(b_num, denom)}"
        return question, pretty_fraction(answer)

    elif op_type == "add_fraction_different_denominator":
        if a==b:
            b += 1
        denoms = range(min(a,b), max(a,b)+1)
        # if len(denoms) == 1:
        #     denom = denoms[0]
        # else:
        #     denom = random.choice(denoms)
        denom_af, denom_bf = random.sample(denoms, 2)  # ensure different denominators
        af_num = random.randint(1, denom_af - 1)
        bf_num = random.randint(1, denom_bf - 1)
        
        af = Fraction(af_num, denom_af)
        bf = Fraction(bf_num, denom_bf)
        answer = af + bf

        question = f"{pretty_fraction(af)} + {pretty_fraction(bf)}"
        return question, pretty_fraction(answer)

    elif op_type == "place_value":
        place_names = ["ones", "tens", "hundreds", "thousands", "ten-thousands", "hundred-thousands", "millions", "ten-millions", "hundred-millions", "billions"]
        digits = list(str(a))
        output = []
        length = len(digits)

        for i, digit in enumerate(digits):
            if digit != '0':
                place_index = length - i - 1
                place = place_names[place_index] if place_index < len(place_names) else f"10^{place_index}"
                output.append(f"{digit} {place}")

        question = f"<font size=75% color=dark_gray>What is</font>\n{' + '.join(output)}?"
        answer = a 
        return question, f"{answer:,}"

    elif op_type == "place_value_reverse":
        place_names = ["ones", "tens", "hundreds", "thousands", "ten-thousands", "hundred-thousands", 
                    "millions", "ten-millions", "hundred-millions", "billions"]

        # Create a number with non-repeating digits
        while True:
            num_digits = len(str(a)) # Use the length of the number generated through the common code as the number of digits
            digits = random.sample("123456789", 1) + random.sample("0123456789", num_digits - 1)
            random.shuffle(digits)
            if len(set(digits)) == len(digits):  # Ensure uniqueness
                break

        number = int("".join(digits))
        digit_list = list(str(number))
        length = len(digit_list)

        # Randomly pick a digit to ask about
        index = random.randint(0, length - 1)
        digit = int(digit_list[index])
        place_index = length - index - 1

        place = place_names[place_index] if place_index < len(place_names) else f"10^{place_index}"
        place_value = digit * (10 ** place_index)

        question = f"<font size=75% color=dark_gray>What value is the </font>\n{digit} in {number:,}?"
        answer = place_value
        return question, f"{answer:,}"
    
    elif op_type == "fraction":
        # Generate fraction addition
        denom1 = generate_number(1, max_val // 2, tier, tier_range) or 1
        denom2 = generate_number(1, max_val // 2, tier, tier_range) or 1
        num1 = generate_number(0, max_val // 2, tier, tier_range)
        num2 = generate_number(0, max_val // 2, tier, tier_range)
        question = f"{num1}/{denom1} + {num2}/{denom2}"
        # Compute answer as simplified fraction
        frac = Fraction(num1, denom1) + Fraction(num2, denom2)
        answer = simplify_fraction(frac.numerator, frac.denominator)
        return question, answer
    return None, None

def generate_question_set(question_type, level):
    """Generate ordered question/answer pairs for a given level."""
    operations = question_type.get(level, {})
    if not operations:
        return []

    # Determine the maximum tier across all operations
    max_tier = max(max(config["tiers"]) for op_configs in operations.values() for config in op_configs)
    questions_by_tier = [[] for _ in range(max_tier + 1)]
    seen_questions = {}  # Dictionary to track questions by op_type, config index, and tier

    for op_type, configs in operations.items():
        # Initialize seen_questions for this op_type if not present
        if op_type not in seen_questions:
            seen_questions[op_type] = {}
        for config_idx, config in enumerate(configs):
            qty = config["qty"]
            min_val = config["min"]
            max_val = config["max"]
            tiers = config["tiers"]
            fontsize = config.get("fontsize", 120)  # Default font size if not specified
            total_tiers = len(tiers)
            if total_tiers == 0:
                continue

            # Initialize seen_questions for this config_idx and tiers
            seen_questions[op_type][config_idx] = {tier: set() for tier in tiers}

            # Distribute questions across tiers
            questions_per_tier = qty // total_tiers
            extra_questions = qty % total_tiers

            for tier in tiers:
                num_questions = questions_per_tier + (1 if tier - tiers[0] < extra_questions else 0)
                for _ in range(num_questions):
                    max_attempts = 100
                    for attempt in range(max_attempts):
                        question, answer = generate_question(op_type, min_val, max_val, tier, tiers)
                        if question and answer:
                            question_key = json.dumps(question, sort_keys=True) if isinstance(question, dict) else question
                            # Check if question is unique for this op_type, config, and tier
                            if question_key not in seen_questions[op_type][config_idx][tier]:
                                seen_questions[op_type][config_idx][tier].add(question_key)
                                questions_by_tier[tier].append((question, answer, tier, fontsize))
                                break
                            # If duplicate, try again
                            if attempt == max_attempts - 1:
                                print(f"Warning: Max attempts reached for {op_type} config {config_idx} in tier {tier}; accepting duplicate: {question}")
                                questions_by_tier[tier].append((question, answer, tier, fontsize))
                                break

    # Shuffle questions within each tier and combine
    result = []
    for tier in range(1, max_tier + 1):
        random.shuffle(questions_by_tier[tier])
        result.extend(questions_by_tier[tier])

    return result

def set_slide_background(slide, red, green, blue):
    # Set background color
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(red, green, blue)

def main():
    # Generate questions for each level
    for level in question_type:
        if int(level) not in levels_to_generate:
            continue

        print(f"\nLevel {level} Questions:")
        
        prev_tier = 0
        question_number = 0
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # Set margins and textbox dimensions for the Question
        left = Inches(0.2)
        top = Inches(0.5)
        width = prs.slide_width - Inches(0.4)  # Leave small margin on sides
        height = Inches(3.5)  # Adjust height as needed

        # --- Year Level Slide ---
        slide_q = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        q_box = add_textbox(slide_q, f"YEAR {level}", left=left, top=Inches(2), width=width, height=Inches(2), default_font_size=120)
        q_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)  # Black text

        questions = generate_question_set(question_type, level)
        for i, (question, answer, tier, fontsize) in enumerate(questions, 1):
            sys.stdout.buffer.write((f"T{tier}Q{i}: {format_question_summary(question)} = {answer}\n").encode('utf-8', errors='replace'))
            if tier != prev_tier:
                tier_count = sum(1 for _, _, t, _ in questions if t == tier)
                # --- Create a Tier slide ---
                slide_t = prs.slides.add_slide(prs.slide_layouts[6])  
                set_slide_background(slide_t, 250, 210, 180)
                add_textbox(slide_t, f"<font size=40%>YEAR {level}</font>\nRound {tier}\n\n<font size=60%>{tier_count} Questions</font>", left=left, top=Inches(1.5), width=width, height=Inches(2), default_font_size=100)
                prev_tier = tier
                question_number = 0

            question_number += 1
            # --- Question Slide ---
            slide_q = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            # Footer text
            add_textbox(slide_q, f"QUESTION {question_number} of {tier_count}", Inches(0.2), prs.slide_height-Inches(1), Inches(3), Inches(1), default_font_size=18, align=PP_ALIGN.LEFT, default_color=COLOR_MAP['dark_gray'])
            # Question Text ---
            if isinstance(question, dict) and question.get("type") == "area":
                q_box = add_textbox(slide_q, "Find the area of the shape. Use the numbers shown.", left, Inches(0.6), width, Inches(1.0), default_font_size=38)
                draw_area_shape(slide_q, question, (prs.slide_width - Inches(7)) / 2, Inches(1.8), Inches(7), Inches(3))
            elif isinstance(question, dict) and question.get("type") == "triangle_angle":
                q_box = add_textbox(slide_q, "Find the missing angle in the triangle.", left, Inches(0.6), width, Inches(1.0), default_font_size=38)
                draw_triangle_angle(slide_q, question, (prs.slide_width - Inches(7)) / 2, (prs.slide_height - Inches(3)) / 2 + Inches(0.5), Inches(7), Inches(3), show_answer=False)
            else:
                q_box = add_textbox(slide_q, question, left, top, width, height, default_font_size=fontsize)
            # Tier information
            t_box = add_textbox(slide_q, f"Round {tier}", prs.slide_width-Inches(3), prs.slide_height-Inches(1), Inches(3), Inches(1), default_font_size=18, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.BOTTOM)
            t_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(180, 180, 180)  # Gray text
            # Add an image (adjust path and size as needed)
            img_path = 'spin_clock.webp'  # e.g. 'my_image.png'
            if os.path.exists(img_path):
                slide_q.shapes.add_picture(img_path, left=width-Inches(0.7), top=Inches(0.1), height=Inches(1))

            # --- Answer Slide ---
            slide_a = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            # Set background color
            set_slide_background(slide_a, 155, 255, 200)
            # Show the question in small text at the top
            clean_question = format_question_summary(question)
            add_textbox(slide_a, f"Q. {clean_question}", left=left, top=Inches(0.3), width=width, height=Inches(1), default_font_size=38)
            if isinstance(question, dict) and question.get("type") == "area":
                draw_area_shape(slide_a, question, (prs.slide_width - Inches(7)) / 2, Inches(1.4), Inches(7), Inches(3))
                a_box = add_textbox(slide_a, answer, left, Inches(4.7), width, Inches(0.7), default_font_size=80)
            elif isinstance(question, dict) and question.get("type") == "triangle_angle":
                # Show triangle with the missing angle revealed and show numeric answer
                draw_triangle_angle(slide_a, question, (prs.slide_width - Inches(7)) / 2, (prs.slide_height - Inches(3)) / 2 + Inches(0.5), Inches(7), Inches(3), show_answer=True)
                a_box = add_textbox(slide_a, answer, left, Inches(4.7), width, Inches(0.7), default_font_size=80)
            else:
                a_box = add_textbox(slide_a, answer, left=left, top=Inches(2.5), width=width, height=Inches(2), default_font_size=120)
            a_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 176, 80)  # Green text
            # Show the question #
            add_textbox(slide_a, f"Answer {question_number} of {tier_count}", 
                        Inches(0.2), prs.slide_height-Inches(1), Inches(3), Inches(1), 
                        default_font_size=18, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.BOTTOM, default_color=COLOR_MAP['dark_gray'])
            # Show the tier
            t_box = add_textbox(slide_a, f"Round {tier}", 
                                prs.slide_width-Inches(3), prs.slide_height-Inches(1), Inches(3), Inches(1), 
                                default_font_size=18, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.BOTTOM, default_color=COLOR_MAP['dark_gray'])
            # t_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(180, 180, 180)  # Gray text

        # Save presentation
        pptx_path = f"c:/temp/data/Maths_Competition_Yr{level}.pptx"
        prs.save(pptx_path)


if __name__ == "__main__":
    random.seed()  # Initialize random seed for reproducibility
    main()